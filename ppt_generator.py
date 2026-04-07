from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import re
import os

class PPTGenerator:
    def __init__(self, theme="modern_blue"):
        self.ppt = Presentation()
        self.title_slide_layout = self.ppt.slide_layouts[0]
        self.title_content_layout = self.ppt.slide_layouts[1]
        self.section_layout = self.ppt.slide_layouts[2] if len(self.ppt.slide_layouts) > 2 else self.ppt.slide_layouts[1]
        self.theme = theme
        self.theme_colors = self._get_theme_colors(theme)
        self.MAX_BULLETS_PER_SLIDE = 7  # Maximum number of bullet points per slide
        
    def _get_theme_colors(self, theme_name):
        """Define color schemes for different themes"""
        themes = {
            "modern_blue": {
                "primary": RGBColor(0, 114, 198),    # Blue
                "secondary": RGBColor(0, 51, 102),   # Dark Blue
                "accent": RGBColor(255, 153, 0),     # Orange
                "background": RGBColor(240, 244, 249),  # Light Blue-Gray
                "text": RGBColor(30, 30, 30)         # Dark Gray
            },
            "elegant_dark": {
                "primary": RGBColor(40, 40, 40),     # Dark Gray
                "secondary": RGBColor(20, 20, 20),   # Almost Black
                "accent": RGBColor(255, 195, 0),     # Gold
                "background": RGBColor(245, 245, 245),  # Almost White
                "text": RGBColor(60, 60, 60)         # Gray
            },
            "vibrant": {
                "primary": RGBColor(213, 0, 82),     # Magenta
                "secondary": RGBColor(35, 35, 35),   # Dark Gray
                "accent": RGBColor(64, 224, 208),    # Turquoise
                "background": RGBColor(248, 248, 248),  # Off-White
                "text": RGBColor(30, 30, 30)         # Almost Black
            },
            "minimal": {
                "primary": RGBColor(70, 70, 70),     # Dark Gray
                "secondary": RGBColor(160, 160, 160), # Medium Gray
                "accent": RGBColor(255, 103, 77),    # Coral
                "background": RGBColor(250, 250, 250), # White
                "text": RGBColor(40, 40, 40)         # Very Dark Gray
            }
        }
        
        return themes.get(theme_name, themes["modern_blue"])
        
    def _process_text_formatting(self, text):
        """Process markdown-style formatting in text"""
        # Bold formatting: **text** -> bold text
        text = re.sub(r'\*\*(.*?)\*\*', r'\1', text)
        
        # Italic formatting: *text* -> italic text
        text = re.sub(r'\*(.*?)\*', r'\1', text)
        
        # Remove any other markdown-style formatting
        text = re.sub(r'__(.*?)__', r'\1', text)
        text = re.sub(r'~~(.*?)~~', r'\1', text)
        
        return text
    
    def _apply_text_formatting(self, paragraph, text):
        """Apply rich text formatting based on markdown-style indicators"""
        # Clear existing text
        for run in paragraph.runs:
            run.text = ""
            
        # Process the basic text without formatting marks
        clean_text = self._process_text_formatting(text)
        paragraph.text = clean_text
        
        # Apply bold formatting
        bold_pattern = r'\*\*(.*?)\*\*'
        bold_matches = re.finditer(bold_pattern, text)
        
        for match in bold_matches:
            for run in paragraph.runs:
                if match.group(1) in run.text:
                    run.font.bold = True
        
        # Apply italic formatting
        italic_pattern = r'\*(.*?)\*'
        italic_matches = re.finditer(italic_pattern, text)
        
        for match in italic_matches:
            for run in paragraph.runs:
                if match.group(1) in run.text:
                    run.font.italic = True
                    
        return paragraph
        
    def _estimate_text_length(self, text):
        """Estimate if text is too long for a single bullet point"""
        # Simple estimation based on character count
        # A typical slide can fit around 80-100 characters per line and about 10 lines
        return len(text) > 100
    
    def _split_long_bullet(self, text):
        """Split a long bullet point into multiple shorter ones at sentence boundaries"""
        sentences = re.split(r'(?<=[.!?])\s+', text)
        if len(sentences) == 1:  # If no sentence boundaries, split on commas or semicolons
            sentences = re.split(r'(?<=[,;])\s+', text)
            
        # If still just one piece, and it's long, try to break at logical points
        if len(sentences) == 1 and len(text) > 100:
            words = text.split()
            sentences = []
            current = []
            char_count = 0
            
            for word in words:
                if char_count + len(word) > 80:  # Aim for ~80 chars per bullet
                    sentences.append(' '.join(current))
                    current = [word]
                    char_count = len(word)
                else:
                    current.append(word)
                    char_count += len(word) + 1  # +1 for the space
                    
            if current:
                sentences.append(' '.join(current))
        
        return sentences
        
    def add_title_slide(self, title, subtitle=None):
        """Add a visually enhanced title slide"""
        slide = self.ppt.slides.add_slide(self.title_slide_layout)
        
        # Add a background shape for visual interest
        left = top = 0
        width = Inches(10)
        height = Inches(0.85)
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.theme_colors["primary"]
        shape.line.color.rgb = self.theme_colors["primary"]
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(44)
        title_shape.text_frame.paragraphs[0].font.color.rgb = self.theme_colors["text"]
        title_shape.text_frame.paragraphs[0].font.bold = True
        
        # Enable text wrapping for title
        title_shape.text_frame.word_wrap = True
        title_shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        
        # Set subtitle if provided
        if subtitle:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = subtitle
            subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)
            subtitle_shape.text_frame.paragraphs[0].font.color.rgb = self.theme_colors["secondary"]
            
            # Enable text wrapping for subtitle
            subtitle_shape.text_frame.word_wrap = True
            subtitle_shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            
        # Add presenter name and date at the bottom
        left = Inches(0.5)
        top = Inches(6.5)
        width = Inches(9)
        height = Inches(0.5)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        p = tf.add_paragraph()
        p.text = "Created with AI Presentation Generator"
        p.alignment = PP_ALIGN.RIGHT
        p.font.size = Pt(12)
        p.font.color.rgb = self.theme_colors["secondary"]

      
        return slide
    
    
    def add_section_slide(self, title, content, image_path=None, slide_number=1, total_slides=1):
        """Add a slide for a section with visually enhanced bullet points and image."""
        slide = self.ppt.slides.add_slide(self.title_content_layout)
        
        # Add a small accent bar at the top
        left = 0
        top = 0
        width = Inches(10)
        height = Inches(0.2)
        accent_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        accent_shape.fill.solid()
        accent_shape.fill.fore_color.rgb = self.theme_colors["accent"]
        accent_shape.line.color.rgb = self.theme_colors["accent"]
        
        # Adjust title if multiple slides for the same section
        display_title = title
        if total_slides > 1:
            left = Inches(9)  # Right side
            top = Inches(6.5)  # Bottom
            width = Inches(0.5)
            height = Inches(0.3)
            
            textbox = slide.shapes.add_textbox(left, top, width, height)
            tf = textbox.text_frame
            p = tf.add_paragraph()
            p.text = f"{slide_number}/{total_slides}"
            p.alignment = PP_ALIGN.RIGHT
            p.font.size = Pt(10)  # Smaller font
            p.font.color.rgb = self.theme_colors["secondary"]
        
        # Set slide title with enhanced styling
        title_shape = slide.shapes.title
        title_shape.text = display_title
        title_shape.text_frame.paragraphs[0].font.size = Pt(36)
        title_shape.text_frame.paragraphs[0].font.color.rgb = self.theme_colors["primary"]
        title_shape.text_frame.paragraphs[0].font.bold = True
        
        # Enable text wrapping for title
        title_shape.text_frame.word_wrap = True
        title_shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        
        # Define dimensions clearly for text and image areas
        slide_width = self.ppt.slide_width
        slide_height = self.ppt.slide_height
        
        margin = Inches(0.5)

         # Determine if this slide should contain an image (starting from slide 3: 3,5,7...)
        has_image = (slide_number >= 3) and (slide_number % 2 == 1) and image_path and os.path.exists(image_path)

        # # Text area dimensions (left side)
        text_width = (slide_width * 2 / 3) - (margin * 1.5)
        text_left = margin
        text_top = margin
        text_height = slide_height - (margin * 2)

        # # Image area dimensions (right side)
        image_width = (slide_width / 3) - (margin * 1.5)
        image_height = image_width * (16 /9)  # Maintain aspect ratio of 16:9
        img_left = text_width + (margin * 2)
        img_top = (slide_height - image_height) / 2  # Vertically center the image

        if has_image or (image_path and os.path.exists(image_path)):
        # Determine positions based on alternating logic starting from slide 3 (right side first)
            
            image_slide_index = ((slide_number - 1) // 2)
            image_slide_index+=1
            if image_slide_index % 2 == 0:
                # Image on right side
                text_left = margin
                img_left = text_width + (margin * 2)
            else:
                # Image on left side
                img_left = margin
                text_left = image_width + (margin * 2)

            img_top = (slide_height - image_height) / 2

        else:
            # No image: Text occupies full width minus margins
            text_left = margin
            text_width = slide_width - (margin * 2)

        # Add Content below Title
        content_top = text_top + Inches(1.5)
        content_box = slide.shapes.add_textbox(text_left, content_top, text_width, text_height - Inches(1.5))
        content_frame = content_box.text_frame  # Correctly reference content frame here!
        content_frame.word_wrap = True

        for idx, point in enumerate(content):
            p = content_frame.add_paragraph()  # Use content_frame instead of text_frame!
            
            clean_point = self._process_text_formatting(point)
            p.text = clean_point
            
            p.level = 0
            p.font.size = Pt(20)
            p.font.color.rgb = self.theme_colors["text"]
            
            self._apply_text_formatting(p, point)
        
        # Add image clearly on right side (1/3 width), vertically centered (2/3 height)
        if image_path and os.path.exists(image_path):
            try:
                pic = slide.shapes.add_picture(image_path, img_left, img_top,
                                            width=image_width,
                                            height=image_height)

            except Exception as e:
                print(f"Error adding image: {str(e)}")
        
        return slide


    
    def add_section_header_slide(self, section_title):
        """Add a divider slide to mark a new section"""
        slide = self.ppt.slides.add_slide(self.section_layout)
        
        # Create a full slide colored background
        left = 0
        top = 0
        width = self.ppt.slide_width
        height = self.ppt.slide_height
        background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        background.fill.solid()
        background.fill.fore_color.rgb = self.theme_colors["primary"]
        background.line.color.rgb = self.theme_colors["primary"]
        
        # Add section title in the center
        left = Inches(1)
        top = Inches(2.5)
        width = Inches(8)
        height = Inches(2)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        tf.word_wrap = True
        
        p = tf.add_paragraph()
        p.text = section_title
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)  # White text
        
        return slide
    
    def add_closing_slide(self, title="Thank You", content=None):
        """Add a visually distinct closing slide"""
        slide = self.ppt.slides.add_slide(self.title_content_layout)
        
        # Set background
        left = 0
        top = 0
        width = self.ppt.slide_width
        height = self.ppt.slide_height
        
        # Add a gradient-like effect with two shapes
        shape1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height/2)
        shape1.fill.solid()
        shape1.fill.fore_color.rgb = self.theme_colors["secondary"]
        shape1.line.color.rgb = self.theme_colors["secondary"]
        
        shape2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, height/2, width, height/2)
        shape2.fill.solid()
        shape2.fill.fore_color.rgb = self.theme_colors["primary"]
        shape2.line.color.rgb = self.theme_colors["primary"]
        
        # Add title in the center
        left = Inches(1)
        top = Inches(2.5)
        width = Inches(8)
        height = Inches(1.5)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.word_wrap = True
        
        p = tf.add_paragraph()
        p.text = title
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(60)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)  # White text
        
        # Add content if provided
        if content:
            left = Inches(1)
            top = Inches(4)
            width = Inches(8)
            height = Inches(1)
            
            content_box = slide.shapes.add_textbox(left, top, width, height)
            tf = content_box.text_frame
            tf.word_wrap = True
            
            p = tf.add_paragraph()
            p.text = content
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(28)
            p.font.color.rgb = RGBColor(255, 255, 255)  # White text
        
        return slide
    
    def _distribute_content(self, title, content, max_slides=None):
        """
        Distribute content across multiple slides if needed, respecting the maximum slides limit.
        
        Args:
            title (str): The section title
            content (list): List of bullet points
            max_slides (int, optional): Maximum number of slides for this section
            
        Returns:
            list: List of tuples (title, content) for each slide
        """
        # Process content to split long bullets
        processed_content = []
        
        for point in content:
            if self._estimate_text_length(point):
                split_points = self._split_long_bullet(point)
                processed_content.extend(split_points)
            else:
                processed_content.append(point)
        
        # If max_slides is specified, we need to adjust content density
        if max_slides and max_slides > 0:
            # Calculate minimum points per slide needed to fit within max_slides
            min_points_per_slide = max(1, len(processed_content) // max_slides)
            # Adjust MAX_BULLETS_PER_SLIDE if needed to fit content within max_slides
            points_per_slide = min(self.MAX_BULLETS_PER_SLIDE, 
                                max(min_points_per_slide, 
                                    min(7, (len(processed_content) + max_slides - 1) // max_slides)))
        else:
            points_per_slide = self.MAX_BULLETS_PER_SLIDE
        
        # Check if we need multiple slides
        if len(processed_content) <= points_per_slide:
            # If content fits on one slide, return it as is
            return [(title, processed_content)]
        
        # Distribute content across multiple slides
        slides_content = []
        num_slides = min(max_slides or float('inf'), 
                        (len(processed_content) + points_per_slide - 1) // points_per_slide)
        
        # Recalculate points per slide to evenly distribute
        points_per_slide = (len(processed_content) + num_slides - 1) // num_slides
        
        for i in range(num_slides):
            start_idx = i * points_per_slide
            end_idx = min((i + 1) * points_per_slide, len(processed_content))
            slide_content = processed_content[start_idx:end_idx]
            slides_content.append((title, slide_content))
        
        return slides_content
    
    def generate_from_content(self, content, images=None):
        """Generate a complete PowerPoint from structured content without section title slides."""
        # Get target slide count
        target_slides = int(content.get("target_slides", 15))

        # Initialize images dictionary if None
        if images is None:
            images = {}

        # Add title slide
        self.add_title_slide(content.get("title", "Presentation"), content.get("subtitle", ""))

        # Calculate fixed slides (title and closing slides only, no section headers)
        fixed_slides = 2  # Title and closing slides

        # Identify sections
        sections = content.get("sections", [])

        # Calculate slides available for content
        available_slides = max(1, target_slides - fixed_slides)

        # Distribute available slides proportionally among sections
        section_weights = [len(section.get("content", [])) for section in sections]
        total_weight = sum(section_weights) or 1  # Avoid division by zero

        # Calculate target slides per section
        section_slides = []
        remaining_slides = available_slides

        for weight in section_weights:
            slides = max(1, int(round((weight / total_weight) * available_slides)))
            if slides > remaining_slides:
                slides = remaining_slides
            section_slides.append(slides)
            remaining_slides -= slides

        # Adjust allocation if necessary
        while remaining_slides < 0:
            idx = section_weights.index(max(section_weights))
            if section_slides[idx] > 1:
                section_slides[idx] -= 1
                remaining_slides += 1
            section_weights[idx] = 0

        while remaining_slides > 0:
            if max(section_weights) > 0:
                idx = section_weights.index(max(section_weights))
                section_slides[idx] += 1
                remaining_slides -= 1
                section_weights[idx] = 0
            else:
                break

        # Create all slides without adding any section header slides
        for idx, section in enumerate(sections):
            section_title = section.get("title", "Section")
            section_content = section.get("content", [])

            distributed_content = self._distribute_content(
                section_title,
                section_content,
                max_slides=section_slides[idx]
            )

            total_section_slides = len(distributed_content)
            for slide_idx, (slide_title, slide_content) in enumerate(distributed_content):
                image_path = images.get(section_title)

                self.add_section_slide(
                    slide_title,
                    slide_content,
                    image_path=image_path,
                    slide_number=slide_idx + 1,
                    total_slides=total_section_slides
                )

        # Add a closing slide with call to action if present
        call_to_action = content.get("call_to_action", "")
        if call_to_action:
            self.add_closing_slide("Thank You", call_to_action)
        else:
            self.add_closing_slide()

        # Return the presentation and actual slide count
        return self.ppt, len(self.ppt.slides)


    def save(self, file_path):
        """Save the generated presentation to a file."""
        self.ppt.save(file_path)